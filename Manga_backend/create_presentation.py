"""
Script to convert the Auto-Manga presentation to PowerPoint format
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return slide

def add_content_slide(title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
    
    return slide

def add_two_column_slide(title, left_content, right_content):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Add left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
    
    # Add right column
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
    
    return slide

# Slide 1: Title
add_title_slide(
    "AUTO-MANGA 🎨",
    "AI-Powered Manga Generation System\n\nLangGraph • Google Gemini • Stable Diffusion • FastAPI • React"
)

# Slide 2: Problem Statement
add_content_slide(
    "Problem Statement & Motivation",
    [
        "⏱️ Manual manga creation: Days to weeks per page",
        "🎨 Requires professional artistic expertise",
        "💰 High production costs",
        "🔄 Difficult to iterate and revise",
        "📐 Hard to maintain character consistency",
        "",
        "✨ Solution: Automated AI-powered manga generation",
        "• Generate one-page manga from text in minutes",
        "• Maintain character consistency across panels",
        "• User-friendly editing and refinement"
    ]
)

# Slide 3: System Architecture
add_content_slide(
    "System Architecture Overview",
    [
        "User Story Input → LangGraph Workflow",
        "",
        "LangGraph Pipeline (8 Nodes):",
        "  1. Story Refinement",
        "  2. Human Review (Interrupt)",
        "  3. Feature Extraction",
        "  4. Character Design",
        "  5. Scene Creation",
        "  6. Director Scripting",
        "  7. Image Prompt Generation",
        "  8. Quality Analysis",
        "",
        "FastAPI Backend → Diffusion Server → Images"
    ]
)

# Slide 4: LangGraph Workflow - Part 1
add_two_column_slide(
    "LangGraph Workflow Pipeline",
    [
        "Node 1: Story Refinement",
        "• Transform user text to manga style",
        "• Add dramatic emotions & onomatopoeia",
        "• Output: 4-6 sentence refined story",
        "",
        "Node 2: Human Review",
        "• Interrupt workflow for user feedback",
        "• Persistent checkpointing",
        "• Resume from checkpoint after edit",
        "",
        "Node 3: Feature Extraction",
        "• Extract story elements",
        "• Characters, setting, conflict",
        "• Mood, tone, sound effects"
    ],
    [
        "Node 4: Character Design",
        "• Detailed character profiles",
        "• Consistency tokens (IBAD_001)",
        "• Physical traits, clothing, expressions",
        "",
        "Node 5: Scene Breakdown",
        "• Break story into 4-5 scenes",
        "• Actions, emotions, dialogues",
        "• Visual elements per scene",
        "",
        "Node 6: Director Script",
        "• Create panel layout (4-5 panels)",
        "• Panel descriptions & camera angles",
        "• Character actions & dialogues"
    ]
)

# Slide 5: LangGraph Workflow - Part 2
add_content_slide(
    "Image Generation & Quality Control",
    [
        "Node 7: Image Prompt Generation",
        "• Create detailed image prompts (≤70 tokens)",
        "• Include consistency tokens",
        "• Manga style keywords (screentone, line art)",
        "• Camera angles, emotions, sound effects",
        "",
        "Node 8: Quality Analysis (Optional)",
        "• Collective page analysis",
        "• Check character consistency across panels",
        "• Evaluate visual continuity & story flow",
        "• Auto-regenerate if score < 7.5/10",
        "• Max 5 regeneration attempts"
    ]
)

# Slide 6: FastAPI Backend
add_content_slide(
    "FastAPI Backend - REST API",
    [
        "POST /api/generate",
        "  • Initial story generation",
        "  • Creates session with unique ID",
        "  • Returns refined story for review",
        "",
        "POST /api/approve",
        "  • Resume workflow with user edits",
        "  • Generate all panel prompts",
        "  • Create images via diffusion server",
        "  • Return image URLs",
        "",
        "POST /api/regenerate-panel",
        "  • Regenerate specific panel with refinements",
        "  • Auto-versioning (panel_1.png → panel_1_1.png)"
    ]
)

# Slide 7: Pydantic Models
add_two_column_slide(
    "Structured Output with Pydantic",
    [
        "MangaFeatureSchema",
        "• main_characters, descriptions",
        "• setting, conflict_or_goal",
        "• mood_and_tone, sound_effects",
        "",
        "CharacterProfile",
        "• Physical traits (face, hair, eyes)",
        "• Clothing, accessories, colors",
        "• consistency_token (unique ID)",
        "• visual_reference_prompt",
        "",
        "SceneFeature",
        "• scene_number, summary",
        "• characters, actions, emotions",
        "• dialogues, inner_thoughts"
    ],
    [
        "Director_Panel",
        "• panel_number, scene_description",
        "• characters_present, actions",
        "• dialogues, sound_effects",
        "• visual_elements",
        "",
        "MangaImagePrompt",
        "• panel_number",
        "• image_prompt (complete description)",
        "",
        "PromptAnalysisResult",
        "• individual panel analyses",
        "• overall_score (1-10)",
        "• needs_regeneration (bool)",
        "• suggested_fixes"
    ]
)

# Slide 8: Character Consistency
add_content_slide(
    "Character Consistency Strategy",
    [
        "Consistency Tokens:",
        "• Each character gets unique ID (IBAD_001, AISHA_002)",
        "• Used in EVERY image prompt",
        "• Prevents character drift across panels",
        "",
        "Visual Reference Prompts:",
        "• Detailed physical description",
        "• Included with consistency token",
        "• Example: 'IBAD_001 (shy anime boy, black hair, navy uniform)'",
        "",
        "Prompt Engineering:",
        "• Same background description across panels",
        "• Consistent lighting and atmosphere",
        "• Manga style keywords in every prompt",
        "• Self-contained panel descriptions"
    ]
)

# Slide 9: Session Management
add_content_slide(
    "Session Management & State Persistence",
    [
        "In-Memory Session Storage:",
        "• Unique session ID per user",
        "• Stores workflow instance & config",
        "• Tracks all intermediate states",
        "",
        "LangGraph Checkpointing:",
        "• InMemorySaver for state persistence",
        "• Resume workflow from any node",
        "• Enables human-in-the-loop via interrupt()",
        "",
        "Session Lifecycle:",
        "1. Create: POST /api/generate",
        "2. Resume: POST /api/approve",
        "3. Regenerate: POST /api/regenerate-panel",
        "4. Cleanup: DELETE /api/session/{id}"
    ]
)

# Slide 10: Quality Analysis
add_content_slide(
    "Quality Analysis & Regeneration",
    [
        "Collective Page Analysis:",
        "• Analyzes ALL panels together (not individually)",
        "• Checks character consistency across panels",
        "• Evaluates visual continuity & story flow",
        "• Assigns quality score 1-10",
        "",
        "Regeneration Logic:",
        "• If overall_score < 7.5 OR 2+ panels need improvement:",
        "  → Regenerate prompts (up to 5 attempts)",
        "• Provides specific feedback for improvement",
        "",
        "Quality Criteria:",
        "✓ Character consistency (same tokens, descriptions)",
        "✓ Visual continuity (backgrounds, lighting)",
        "✓ Manga style elements (screentone, line art)",
        "✓ Logical narrative progression"
    ]
)

# Slide 11: Complete Workflow
add_content_slide(
    "End-to-End Process Flow",
    [
        "1. User submits story → POST /api/generate",
        "2. Story refinement (LLM transforms to manga style)",
        "3. Review checkpoint → Returns to user",
        "4. User edits story → POST /api/approve",
        "5. Workflow resumes from checkpoint",
        "6. Feature extraction → Character design",
        "7. Scene breakdown → Director scripting",
        "8. Image prompt generation",
        "9. Quality analysis (optional regeneration)",
        "10. Generate images via diffusion server",
        "11. Save images to disk",
        "12. Return URLs to frontend",
        "",
        "⏱️ Total time: 2-3 minutes per page"
    ]
)

# Slide 12: Technology Stack
add_two_column_slide(
    "Technology Stack",
    [
        "Backend:",
        "• FastAPI - Web framework",
        "• LangGraph - Workflow orchestration",
        "• LangChain - LLM integration",
        "• Google Gemini 2.5 - Language model",
        "• Pydantic 2.0+ - Data validation",
        "• Python 3.10+",
        "",
        "Frontend:",
        "• React 18+ - UI framework",
        "• Vite - Build tool",
        "• JavaScript/TypeScript",
        "• Tailwind CSS"
    ],
    [
        "Machine Learning:",
        "• Stable Diffusion (remote server)",
        "• PIL (Pillow) - Image handling",
        "",
        "Infrastructure:",
        "• Docker - Containerization",
        "• PostgreSQL (planned)",
        "• Redis (planned)",
        "",
        "Performance:",
        "• LLM: 16-26s, ~3,800 tokens",
        "• Image gen: 60-120s (4 panels)",
        "• Total: ~2-3 minutes/page"
    ]
)

# Slide 13: Key Achievements
add_content_slide(
    "Technical Achievements",
    [
        "✅ Sophisticated Workflow Orchestration",
        "  • 8-node pipeline with checkpointing",
        "  • Human-in-the-loop via REST API",
        "  • Resumable across HTTP requests",
        "",
        "✅ Character Consistency Innovation",
        "  • Consistency tokens prevent drift",
        "  • 85-90% visual consistency across panels",
        "",
        "✅ Quality Assurance System",
        "  • Collective page analysis (not individual)",
        "  • Automatic regeneration with feedback",
        "  • ~20% regeneration rate",
        "",
        "✅ Production-Ready Architecture",
        "  • Complete error handling & CORS",
        "  • Session management & versioning",
        "  • Static file serving"
    ]
)

# Slide 14: Business Value
add_content_slide(
    "Business Value & Impact",
    [
        "Time Reduction: Days → Minutes",
        "• Manual: Days to weeks per page",
        "• Auto-Manga: 2-3 minutes per page",
        "",
        "Cost Savings:",
        "• Eliminates need for professional artist",
        "• ~$0.15 per page (LLM + API costs)",
        "",
        "Accessibility:",
        "• Non-artists can create visual stories",
        "• Instant iteration and refinement",
        "• Democratizes manga creation",
        "",
        "Use Cases:",
        "• Content creators: Fast prototyping",
        "• Publishing: Quick story adaptations",
        "• Education: Visualize narratives",
        "• Entertainment: Fan art generation"
    ]
)

# Slide 15: Comparison
add_content_slide(
    "Comparison with Alternatives",
    [
        "Manual Creation:",
        "  Time: Days | Cost: $$$$ | Consistency: High | Access: Low",
        "",
        "Other AI Tools (e.g., ComicAI):",
        "  Time: Hours | Cost: $$ | Consistency: Medium | Access: Medium",
        "",
        "Auto-Manga:",
        "  Time: Minutes | Cost: $ | Consistency: High | Access: High",
        "",
        "Unique Advantages:",
        "✓ True human-in-the-loop workflow",
        "✓ Consistency tokens for character stability",
        "✓ Collective quality analysis",
        "✓ Panel-level refinement",
        "✓ Production-ready REST API"
    ]
)

# Slide 16: Future Roadmap
add_content_slide(
    "Future Enhancements",
    [
        "Short Term (1-2 months):",
        "• Style selection (manga, comic, anime)",
        "• Database-backed sessions (PostgreSQL)",
        "• User authentication & dashboard",
        "",
        "Medium Term (3-6 months):",
        "• Multi-page story support (5-10 pages)",
        "• Character reference image upload",
        "• Story variation generation",
        "• Export to PDF/ebook formats",
        "",
        "Long Term (6+ months):",
        "• Mobile apps (iOS/Android)",
        "• Community sharing platform",
        "• API marketplace & monetization",
        "• Publishing platform integration"
    ]
)

# Slide 17: Conclusion
add_content_slide(
    "Key Takeaways",
    [
        "Auto-Manga demonstrates:",
        "",
        "🎯 Advanced AI orchestration with LangGraph",
        "🔄 Production-ready human-in-the-loop workflows",
        "🎨 Innovative consistency techniques",
        "📊 Intelligent quality assurance systems",
        "🚀 Full-stack modern architecture",
        "",
        "Learning outcomes:",
        "• Multi-agent systems & LLM orchestration",
        "• Prompt engineering for consistency",
        "• State management in stateless APIs",
        "• Quality analysis for AI outputs",
        "",
        "Let's create manga, one panel at a time! 🎨✨"
    ]
)

# Slide 18: Thank You
slide_layout = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(slide_layout)

# Add centered "Thank You" text
thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.text = "Thank You!\n\nQuestions & Discussion"
thank_you_frame.paragraphs[0].font.size = Pt(54)
thank_you_frame.paragraphs[0].font.bold = True
thank_you_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
thank_you_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# Add contact/footer
footer_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
footer_frame = footer_box.text_frame
footer_frame.text = "Auto-Manga: AI-Powered Manga Generation System\nLangGraph • FastAPI • Google Gemini • Stable Diffusion"
footer_frame.paragraphs[0].font.size = Pt(14)
footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r"c:\Users\Haider\Desktop\Manga\LangGraph\Manga_backend\AUTO_MANGA_PRESENTATION.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
